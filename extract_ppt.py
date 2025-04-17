import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import colorsys
from collections import Counter

def extract_dominant_colors(image_path, num_colors=5):
    img = Image.open(image_path)
    img = img.convert('RGB')
    pixels = list(img.getdata())
    color_counts = Counter(pixels)
    dominant_colors = color_counts.most_common(num_colors)
    return [{'rgb': color, 'count': count} for color, count in dominant_colors]

def rgb_to_hex(rgb):
    return '#{:02x}{:02x}{:02x}'.format(rgb[0], rgb[1], rgb[2])

def process_ppt(ppt_path):
    prs = Presentation(ppt_path)
    output_dir = 'extracted_content'
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)
    
    md_content = []
    all_colors = []
    image_counter = 0
    
    for slide_number, slide in enumerate(prs.slides, 1):
        md_content.append(f'\n## Slide {slide_number}\n')
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        md_content.append(text + '\n')
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_counter += 1
                image_path = os.path.join(images_dir, f'image_{image_counter}.png')
                image = shape.image
                with open(image_path, 'wb') as f:
                    f.write(image.blob)
                md_content.append(f'![Slide {slide_number} Image {image_counter}](images/image_{image_counter}.png)\n')
                
                # 分析圖片顏色
                colors = extract_dominant_colors(image_path)
                all_colors.extend([color['rgb'] for color in colors])
    
    # 寫入Markdown文件
    with open(os.path.join(output_dir, 'content.md'), 'w', encoding='utf-8') as f:
        f.write(''.join(md_content))
    
    # 分析主題色
    color_counter = Counter(all_colors)
    dominant_colors = color_counter.most_common(5)
    
    # 生成配色建議
    with open(os.path.join(output_dir, 'color_theme.md'), 'w', encoding='utf-8') as f:
        f.write('# 主題配色建議\n\n')
        f.write('## 主要顏色\n')
        for i, (color, count) in enumerate(dominant_colors, 1):
            hex_color = rgb_to_hex(color)
            f.write(f'{i}. {hex_color} (使用頻率: {count})\n')
        
        f.write('\n## UI/UX 建議\n')
        f.write('1. 主色調：使用最常見的顏色作為主要品牌色\n')
        f.write('2. 輔助色：使用第二和第三常見的顏色作為輔助色\n')
        f.write('3. 強調色：使用較少出現的鮮艷顏色作為強調色\n')
        f.write('4. 背景色：建議使用主色調的淺色版本作為背景\n')
        f.write('5. 文字色：根據背景色的明暗度選擇適當的文字顏色\n')

if __name__ == '__main__':
    ppt_path = 'Jeff.pptx'
    process_ppt(ppt_path)
    print('提取完成！請查看 extracted_content 目錄下的文件。')