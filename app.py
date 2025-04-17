import os
import io
import zipfile
import logging
import markdown
from flask import Flask, render_template, request, send_file, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from collections import Counter
import shutil
import tempfile
import traceback
from flask_cors import CORS

# 配置日誌
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})  # 啟用CORS以允許Next.js前端連接
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024  # 32MB max upload size
MAX_FILE_SIZE_MB = 30  # 展示給用戶的最大文件大小限制

# 確保上傳和輸出目錄存在
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

# 從原始腳本中提取的函數
def extract_dominant_colors(image_path, num_colors=5):
    try:
        img = Image.open(image_path)
        img = img.convert('RGB')
        pixels = list(img.getdata())
        color_counts = Counter(pixels)
        dominant_colors = color_counts.most_common(num_colors)
        return [{'rgb': color, 'count': count} for color, count in dominant_colors]
    except Exception as e:
        logger.error(f"提取顏色時出錯: {str(e)}")
        return [{'rgb': (200, 200, 200), 'count': 1}]  # 返回默認顏色

def rgb_to_hex(rgb):
    return '#{:02x}{:02x}{:02x}'.format(rgb[0], rgb[1], rgb[2])

def process_ppt(ppt_path, output_dir, output_id):
    prs = Presentation(ppt_path)
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)
    
    md_content = []
    all_colors = []
    image_counter = 0
    
    for slide_number, slide in enumerate(prs.slides, 1):
        md_content.append(f'\n## Slide {slide_number}\n')
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        md_content.append(text + '\n')
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_counter += 1
                    image_path = os.path.join(images_dir, f'image_{image_counter}.png')
                    image = shape.image
                    with open(image_path, 'wb') as f:
                        f.write(image.blob)
                    
                    # 使用相對路徑引用圖片，這樣在ZIP文件中可以正確顯示
                    md_content.append(f'![Slide {slide_number} Image {image_counter}](images/image_{image_counter}.png)\n')
                    
                    # 分析圖片顏色
                    colors = extract_dominant_colors(image_path)
                    all_colors.extend([color['rgb'] for color in colors])
                except Exception as e:
                    logger.error(f"處理圖片時出錯: {str(e)}")
                    md_content.append(f'*圖片處理失敗: Slide {slide_number}*\n')
    
    # 寫入Markdown文件
    with open(os.path.join(output_dir, 'content.md'), 'w', encoding='utf-8') as f:
        f.write(''.join(md_content))
    
    # 分析主題色
    color_counter = Counter(all_colors)
    dominant_colors = color_counter.most_common(5) if all_colors else [(200, 200, 200)]
    
    # 生成配色建議
    with open(os.path.join(output_dir, 'color_theme.md'), 'w', encoding='utf-8') as f:
        f.write('# 主題配色建議\n\n')
        f.write('## 主要顏色\n')
        for i, (color, count) in enumerate(dominant_colors, 1):
            hex_color = rgb_to_hex(color)
            f.write(f'{i}. {hex_color} (使用頻率: {count})\n')
        
        f.write('\n## UI/UX 建議\n')
        f.write('1. 主色調：使用最常見的顏色作為主要品牌色\n')
        f.write('2. 輔助色：使用第二和第三常見的顏色作為輔助色\n')
        f.write('3. 強調色：使用較少出現的鮮艷顏色作為強調色\n')
        f.write('4. 背景色：建議使用主色調的淺色版本作為背景\n')
        f.write('5. 文字色：根據背景色的明暗度選擇適當的文字顏色\n')
    
    return {
        'content_md': os.path.join(output_dir, 'content.md'),
        'color_theme_md': os.path.join(output_dir, 'color_theme.md'),
        'images_dir': images_dir,
        'dominant_colors': [{'hex': rgb_to_hex(color), 'count': count} for color, count in dominant_colors]
    }

@app.route('/')
def index():
    return render_template('index.html')

# 確保支持API前綴的路由
@app.route('/api/<output_id>/images/<filename>')
@app.route('/<output_id>/images/<filename>')
def serve_image(output_id, filename):
    logger.info(f"請求圖片: /{output_id}/images/{filename}")
    output_dir = os.path.join(app.config['OUTPUT_FOLDER'], output_id)
    images_dir = os.path.join(output_dir, 'images')
    if os.path.exists(os.path.join(images_dir, filename)):
        logger.info(f"找到圖片: {os.path.join(images_dir, filename)}")
    else:
        logger.error(f"圖片不存在: {os.path.join(images_dir, filename)}")
    return send_from_directory(images_dir, filename)

# 支持原始/api路徑
@app.route('/api/upload', methods=['POST'])
@app.route('/upload', methods=['POST'])  
def upload_file():
    logger.info("接收到上傳請求")
    if 'file' not in request.files:
        logger.warning("上傳失敗: 沒有選擇文件")
        return jsonify({'error': '沒有選擇文件'}), 400
    
    file = request.files['file']
    if file.filename == '':
        logger.warning("上傳失敗: 沒有選擇文件")
        return jsonify({'error': '沒有選擇文件'}), 400
    
    if not file.filename.endswith('.pptx'):
        logger.warning(f"上傳失敗: 文件類型不支持 - {file.filename}")
        return jsonify({'error': '只支持PPTX文件'}), 400
    
    try:
        # 檢查文件大小
        file_content = file.read()
        file_size_mb = len(file_content) / (1024 * 1024)
        logger.info(f"上傳文件: {file.filename}, 大小: {file_size_mb:.2f}MB")
        
        if file_size_mb > MAX_FILE_SIZE_MB:
            return jsonify({'error': f'文件過大。最大允許大小為{MAX_FILE_SIZE_MB}MB，您的檔案為{file_size_mb:.1f}MB'}), 413
        
        # 重置文件指針以便後續處理
        file.seek(0)
        
        # 創建唯一的輸出目錄
        output_id = str(hash(file.filename + str(os.urandom(8))))[1:10]
        output_dir = os.path.join(app.config['OUTPUT_FOLDER'], output_id)
        os.makedirs(output_dir, exist_ok=True)
        logger.info(f"創建輸出目錄: {output_dir}")
        
        # 保存上傳的文件
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        logger.info(f"文件已保存至: {file_path}")
        
        try:
            # 處理PPT文件
            logger.info(f"開始處理PPT文件: {file_path}")
            result = process_ppt(file_path, output_dir, output_id)
            logger.info(f"PPT處理完成")
            
            # 創建ZIP文件
            zip_path = os.path.join(app.config['OUTPUT_FOLDER'], f'{output_id}.zip')
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                # 添加Markdown文件
                zipf.write(result['content_md'], os.path.basename(result['content_md']))
                zipf.write(result['color_theme_md'], os.path.basename(result['color_theme_md']))
                
                # 添加圖片
                for img_file in os.listdir(result['images_dir']):
                    img_path = os.path.join(result['images_dir'], img_file)
                    zipf.write(img_path, os.path.join('images', img_file))
            logger.info(f"ZIP文件已創建: {zip_path}")
            
            # 讀取Markdown內容用於預覽
            with open(result['content_md'], 'r', encoding='utf-8') as f:
                content_md = f.read()
            
            with open(result['color_theme_md'], 'r', encoding='utf-8') as f:
                color_theme_md = f.read()
            
            # 嘗試使用 markdown 模組，如果不可用則回傳原始文字
            try:
                import markdown
                # 修改圖片路徑，將相對路徑 'images/image_X.png' 轉換為 '/{output_id}/images/image_X.png'
                # 這樣在前端顯示時才能正確加載圖片
                content_md_with_api_paths = content_md.replace('](images/', f'](/{output_id}/images/')
                # 檢查並移除可能存在的 /api/ 前綴，以確保路徑與後端期望一致
                content_md_with_api_paths = content_md_with_api_paths.replace(f'](/api/{output_id}/images/', f'](/{output_id}/images/')
                html_content = markdown.markdown(content_md_with_api_paths)
            except ImportError:
                logger.warning("markdown 模組未安裝，將返回原始文本")
                html_content = content_md
            
            return jsonify({
                'success': True,
                'output_id': output_id,
                'content_md': content_md,
                'html_content': html_content,
                'color_theme_md': color_theme_md,
                'dominant_colors': result['dominant_colors'],
                'original_filename': filename
            })
        
        except Exception as e:
            error_details = traceback.format_exc()
            logger.error(f"處理文件時出錯: {str(e)}\n{error_details}")
            return jsonify({'error': f'處理文件時出錯: {str(e)}'}), 500
        finally:
            # 清理上傳的文件
            if os.path.exists(file_path):
                os.remove(file_path)
                logger.info(f"已刪除原始上傳文件: {file_path}")
    
    except Exception as e:
        error_details = traceback.format_exc()
        logger.error(f"上傳處理過程中出錯: {str(e)}\n{error_details}")
        return jsonify({'error': f'上傳處理過程中出錯: {str(e)}'}), 500

# 支持原始路徑和API前綴路徑
@app.route('/download/<output_id>')
@app.route('/api/download/<output_id>')
def download_file(output_id):
    logger.info(f"請求下載文件: {output_id}")
    zip_path = os.path.join(app.config['OUTPUT_FOLDER'], f'{output_id}.zip')
    if not os.path.exists(zip_path):
        logger.warning(f"下載失敗: 文件不存在 - {zip_path}")
        return jsonify({'error': '文件不存在或已過期'}), 404
    
    # 從請求中取得原始文件名
    filename = request.args.get('filename', 'extracted_content')
    # 添加byaslaninno後綴
    download_name = f"{filename}_byaslaninno.zip"
    logger.info(f"提供下載: {download_name}")
    
    return send_file(zip_path, as_attachment=True, download_name=download_name)

# 清理臨時文件的路由（可選）
@app.route('/cleanup/<output_id>', methods=['POST'])
@app.route('/api/cleanup/<output_id>', methods=['POST'])
def cleanup(output_id):
    logger.info(f"請求清理文件: {output_id}")
    output_dir = os.path.join(app.config['OUTPUT_FOLDER'], output_id)
    zip_path = os.path.join(app.config['OUTPUT_FOLDER'], f'{output_id}.zip')
    
    if os.path.exists(output_dir):
        shutil.rmtree(output_dir)
        logger.info(f"已刪除輸出目錄: {output_dir}")
    
    if os.path.exists(zip_path):
        os.remove(zip_path)
        logger.info(f"已刪除ZIP文件: {zip_path}")
    
    return jsonify({'success': True})

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=8080)